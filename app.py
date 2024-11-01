from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
import pandas as pd
import requests
from io import BytesIO

# פונקציה ליצירת שקופיות במצגת מתוך נתוני Excel
def create_presentation_from_excel(excel_path, output_pptx):
    # קרא את קובץ ה-Excel, תוך התעלמות מהשורות הראשונות (עד שורה 3)
    data = pd.read_excel(excel_path, skiprows=3)
    
    # צור מצגת חדשה
    prs = Presentation()
    
    for index, row in data.iterrows():
        # צור שקופית חדשה
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Layout מספר 5 לרקע ריק

        # כותרת
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = row[0]  # כותרת מהעמודה הראשונה (A)
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # תאריך פרסום
        date_box = slide.shapes.add_textbox(Inches(8), Inches(6.5), Inches(2), Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = f"תאריך פרסום: {row[3]}"  # תאריך מהעמודה הרביעית (D)
        date_frame.paragraphs[0].font.size = Pt(14)
        date_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # כלי תקשורת
        media_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(0.5))
        media_frame = media_box.text_frame
        media_frame.text = f"כלי תקשורת: {row[4]}"  # כלי תקשורת מהעמודה החמישית (E)
        media_frame.paragraphs[0].font.size = Pt(14)
        media_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # תמונה
        image_url = row[9]  # קישור לתמונה מהעמודה העשירית (J)
        if pd.notnull(image_url):  # בדיקה אם יש קישור לתמונה
            try:
                response = requests.get(image_url)
                image_stream = BytesIO(response.content)
                slide.shapes.add_picture(image_stream, Inches(3), Inches(3), Inches(4), Inches(3))  # מרכז את התמונה
            except Exception as e:
                print(f"Error loading image from URL: {image_url} - {e}")
        
        # קישור
        link_text = row[8]  # קישור מהעמודה התשיעית (I)
        if pd.notnull(link_text):  # רק אם יש קישור
            link_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(0.5))
            link_frame = link_box.text_frame
            link_frame.text = f"קישור: {link_text}"
            link_frame.paragraphs[0].font.size = Pt(14)
            link_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # שמור את קובץ ה-PowerPoint
    prs.save(output_pptx)

# דוגמה לשימוש
create_presentation_from_excel("data.xlsx", "output_presentation.pptx")
